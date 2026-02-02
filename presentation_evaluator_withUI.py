# プレゼン評価システム

import streamlit as st
import os
import sys
import re
import pptx
import openai
import numpy as np
import base64
import pandas as pd
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from datetime import datetime

# ==== ページ設定 ====
st.set_page_config(page_title="AIプレゼン評価システム", layout="wide")

# ==== スタイル設定 ====
st.markdown("""
    <style>
    .main {
        background-color: #f8f9fa;
    }
    .stMetric {
        background-color: #ffffff;
        padding: 15px;
        border-radius: 10px;
        box-shadow: 0 2px 4px rgba(0,0,0,0.05);
    }
    </style>
    """, unsafe_allow_html=True)

# ==== タイトル・説明 ====
st.title("🎤 AIプレゼン評価システム")
st.markdown("音声ファイルとスライド資料をアップロードするだけで、AIがあなたのプレゼンを多角的に分析・採点します。")

# ==== サイドバー設定 ====
with st.sidebar:
    st.header("⚙️ 設定")
    api_key = st.text_input("OpenAI API Keyを入力してください", type="password")
    
    # 元のコードのモデル設定を反映（2026年時点の最新を想定）
    model_llm = st.selectbox("使用モデル (LLM)", ["gpt-5.2-2025-12-11", "gpt-5-nano"], index=0)
    model_whisper = "whisper-1"
    
    st.info("""
    **分析項目:**
    1. 内容 (30%)
    2. プレゼン技術 (30%)
    3. 視覚資料 (20%)
    4. 構成 (20%)
    """)

# APIキーのセット
if api_key:
    client = openai.OpenAI(api_key=api_key)
else:
    st.warning("⚠️ 続行するにはサイドバーにOpenAI APIキーを入力してください。")
    st.stop()


# ==== 音声分析モジュール ====
def transcribe_audio(file_path):
    audio_file = open(file_path, "rb")
    response = openai.audio.transcriptions.create(
        model=model_whisper,
        file=audio_file,
        response_format="verbose_json",
        language="ja"
    )

    text = response.text
    segments = response.segments

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"transcription_{timestamp}.txt"
    with open(filename, "w", encoding="utf-8") as f:
        f.write(text)

    return text, segments


def analyze_speech(segments):
    total_words = sum(len(seg.text.split()) for seg in segments)
    duration_minutes = (segments[-1].end - segments[0].start) / 60.0
    wpm = total_words / duration_minutes if duration_minutes else 0

    filler_words = ['えーと', 'あの', 'えっと', 'その']
    filler_count = sum(sum(word in seg.text for word in filler_words) for seg in segments)

    pause_lengths = [segments[i + 1].start - segments[i].end for i in range(len(segments) - 1)]
    long_pauses = sum(1 for p in pause_lengths if p > 1.0)

    return {
        "wpm": round(wpm, 2),
        "filler_count": filler_count,
        "long_pauses": long_pauses
    }


# ==== 資料抽出モジュール ====
def extract_ppt_text(file_path):
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slides_text.append(f"スライド {i + 1}:\n{slide_text}\n")
    return "\n".join(slides_text)


def extract_images_from_ppt(ppt_path, output_dir):
    prs = Presentation(ppt_path)
    image_files = []

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_filename = os.path.join(output_dir, f"slide_{i + 1}_image_{j + 1}.png")
                with open(image_filename, 'wb') as f:
                    f.write(image_bytes)
                image_files.append(image_filename)

    return image_files


def encode_image_to_base64(image_path):
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')


def analyze_image(image_path):
    st.write("ファイル名: {image_path}")
    base64_image = encode_image_to_base64(image_path)
    response = client.chat.completions.create(
        model=model_llm,
        messages=[
            {"role": "system", "content": "あなたは画像解析の専門家です。"},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "この画像に何が写っているか説明し、プレゼン資料として適切か評価してください。視覚資料としての質も100点満点（整数）で採点してください。フォーマット: 視覚資料: ○点"},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                ]
            }
        ]
    )
    return response.choices[0].message.content


def extract_visual_score(image_analysis):
    pattern = r"視覚資料: ([0-9]{1,3})点"
    matches = re.findall(pattern, image_analysis)
    if matches:
        scores = [int(score) for score in matches]
        return int(sum(scores) / len(scores))
    return 0


def analyze_all_images(image_files):
    all_analyses = []
    for image_path in image_files:
        analysis = analyze_image(image_path)
        all_analyses.append(f"{image_path}:\n{analysis}\n")
    return "\n".join(all_analyses)


def analyze_slide_text(slide_text):
    response = client.chat.completions.create(
        model=model_llm,
        messages=[
            {"role": "system", "content": "あなたはプロのプレゼン資料評価者です。"},
            {"role": "user", "content": f"""
以下はプレゼンテーションのスライド全文です。

[スライド全文]
{slide_text}

この資料のスライド数、各スライドの文字量の適切さ、内容を評価し、全体的な資料の質を以下のフォーマットで100点満点（整数）で評価してください。
ただし、図表や画像は評価に含めないでください。

資料: ○点

その後に資料の良い点と改善点を簡単にまとめてください。
"""}
        ]
    )
    return response.choices[0].message.content


def generate_evaluation_with_images(transcription, slide_text_analysis, image_analysis):
    response = client.chat.completions.create(
        model=model_llm,
        messages=[
            {"role": "system", "content": "あなたはプロのプレゼン評価者です。"},
            {"role": "user", "content": f"""
以下はプレゼンの文字起こしとスライド資料の分析結果、および画像分析結果です。

[文字起こし]:
{transcription}

[スライドテキスト分析]:
{slide_text_analysis}

[画像分析]:
{image_analysis}

以下の4つの観点（内容、プレゼン技術、視覚資料、構成）について、それぞれ100点満点（整数）で評価し、簡単な理由と改善点、長所を出力してください。
最後に3つの改善点と具体的なアドバイスも示してください。

フォーマットは必ず以下としてください：
内容: ○点
プレゼン技術: ○点
視覚資料: ○点
構成: ○点

その後に評価コメントを書いてください。
"""}
        ]
    )
    return response.choices[0].message.content


# ==== スコア抽出 ====
def extract_scores(evaluation_text):
    pattern = r"内容: ([0-9]{1,3})点.*?プレゼン技術: ([0-9]{1,3})点.*?視覚資料: ([0-9]{1,3})点.*?構成: ([0-9]{1,3})点"
    match = re.search(pattern, evaluation_text, re.DOTALL)

    if match:
        return {
            "内容": int(match.group(1)),
            "プレゼン技術": int(match.group(2)),
            "視覚資料": int(match.group(3)),
            "構成": int(match.group(4))
        }
    return {"内容": 0, "プレゼン技術": 0, "視覚資料": 0, "構成": 0}


def compute_score(sub_scores):
    weights = {"内容": 0.3, "プレゼン技術": 0.3, "視覚資料": 0.2, "構成": 0.2}
    total = sum(float(sub_scores[k]) * weights[k] for k in weights)
    return int(round(total, 0))

# ==== Web UI メイン処理 ====
col1, col2 = st.columns(2)
with col1:
    audio_upload = st.file_uploader("1. 音声ファイルをアップロード", type=['mp3', 'wav', 'm4a', 'mp4'])
with col2:
    ppt_upload = st.file_uploader("2. PowerPointファイルをアップロード", type=['pptx'])

if st.button("📊 プレゼンを分析する", use_container_width=True):
    if not api_key:
        st.error("APIキーを入力してください")
        st.stop()

    os.environ["OPENAI_API_KEY"] = api_key 
    client = openai.OpenAI(api_key=api_key)

    if audio_upload and ppt_upload:
        # 一時ファイル保存用ディレクトリ
        temp_dir = "temp_process"
        if not os.path.exists(temp_dir): os.makedirs(temp_dir)
        
        audio_path = os.path.join(temp_dir, audio_upload.name)
        ppt_path = os.path.join(temp_dir, ppt_upload.name)
        
        with open(audio_path, "wb") as f: f.write(audio_upload.getbuffer())
        with open(ppt_path, "wb") as f: f.write(ppt_upload.getbuffer())

        try:
            with st.status("分析中...", expanded=True) as status:
                st.write("🎙️ 音声をテキスト化・分析中...")
                text, segments = transcribe_audio(audio_path)
                speech_analysis = analyze_speech(segments)

                st.write("📄 スライドテキストを抽出中...")
                slides_text = extract_ppt_text(ppt_path)
                slide_text_analysis = analyze_slide_text(slides_text)

                st.write("🖼️ 画像を解析中...")
                img_extract_dir = os.path.join(temp_dir, "extracted_images")
                image_files = extract_images_from_ppt(ppt_path, img_extract_dir)
                
                if image_files:
                    image_analysis = analyze_all_images(image_files)
                    image_visual_score = extract_visual_score(image_analysis)
                else:
                    image_analysis = "画像は含まれていません。"
                    image_visual_score = 0

                st.write("🤖 総合評価を生成中...")
                evaluation = generate_evaluation_with_images(text, slide_text_analysis, image_analysis)
                
                # スコア計算
                sub_scores = extract_scores(evaluation)
                if image_visual_score == 0:
                    image_visual_score = sub_scores["視覚資料"]
                sub_scores["視覚資料"] = int(round((sub_scores["視覚資料"] + image_visual_score) / 2, 0))
                total_score = compute_score(sub_scores)

                status.update(label="✅ 分析が完了しました！", state="complete", expanded=False)

            # ==== 結果保存セクション ====
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            result_filename = f"evaluation_result_{timestamp}.txt"

            with open(result_filename, "w", encoding="utf-8") as f:
                f.write(f"==== 総合得点: {total_score}点 ====\n\n")
                f.write("==== 総合評価 ====\n")
                f.write(evaluation + "\n\n")
                f.write("==== 音声分析 ====\n")
                f.write(str(speech_analysis) + "\n\n")
#                f.write("==== スライドテキスト分析 ====\n")    #テキストのみの分析結果とならないため非表示
#                f.write(slide_text_analysis + "\n\n")
#                f.write("==== 画像分析 ====\n")               #冗長な結果しか出力できないため非表示
#                f.write(image_analysis + "\n\n")

            print(f"\n評価結果をファイルに保存しました: {result_filename}")

            # ==== 結果表示セクション ====
            st.divider()
            
            # 総合得点の表示
            c1, c2 = st.columns([1, 2])
            with c1:
                st.metric(label="総合得点", value=f"{total_score} 点")
                # レーダーチャート用データ準備
                score_data = pd.DataFrame({
                    "項目": list(sub_scores.keys()),
                    "得点": list(sub_scores.values())
                })
                st.bar_chart(score_data.set_index("項目"))

            with c2:
                st.subheader("🔊 音声分析")
                sc1, sc2, sc3 = st.columns(3)
                sc1.metric("話す速さ (WPM)", speech_analysis['wpm'])
                sc2.metric("フィラー数", speech_analysis['filler_count'])
                sc3.metric("長い沈黙", speech_analysis['long_pauses'])

            st.divider()
            
            # タブによる詳細表示
            tab1, tab2, tab3 = st.tabs(["📝 総合評価レポート", "📖 文字起こし全文", "🖼️ スライド分析詳細"])
            
            with tab1:
                st.markdown(evaluation)
                
            with tab2:
                st.text_area("文字起こし内容", text, height=300)
                
            with tab3:
                st.markdown("### スライドテキスト評価")
                st.write(slide_text_analysis)
                if image_files:
                    st.markdown("### 抽出された画像とAIコメント")
                    for img in image_files:
                        st.image(img, width=300)

        except Exception as e:
            st.error(f"分析中にエラーが発生しました: {e}")
        
        finally:
            # クリーンアップ
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)

    else:
        st.info("音声ファイルとPowerPointファイルを両方アップロードして、分析ボタンを押してください。")

# ==== フッター ====
st.markdown("---")
st.caption(f"Presentation Evaluator Pro v2.0 | Powered by {model_llm}")
