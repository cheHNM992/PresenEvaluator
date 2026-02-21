# ãƒ—ãƒ¬ã‚¼ãƒ³è©•ä¾¡ã‚·ã‚¹ãƒ†ãƒ  (CLI/GUIçµ±åˆç‰ˆ)
# å¯¾å¿œLLM: OpenAI, OpenRouter, Ollama

import os
import sys
import re
import base64
import shutil
from datetime import datetime
from enum import Enum

# StreamlitãŒã‚¤ãƒ³ãƒãƒ¼ãƒˆå¯èƒ½ã‹ãƒã‚§ãƒƒã‚¯
try:
    import streamlit as st
    STREAMLIT_AVAILABLE = True
except ImportError:
    STREAMLIT_AVAILABLE = False

import openai
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ==== LLMãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ ====
class LLMProvider(Enum):
    OPENAI = "openai"
    OPENROUTER = "openrouter"
    OLLAMA = "ollama"


# ==== ã‚°ãƒ­ãƒ¼ãƒãƒ«è¨­å®š ====
# ä½¿ç”¨ã™ã‚‹LLMãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼è¨­å®š
PROVIDER = "openai"     # openai, openrouter, ollama

# OpenAIç”¨è¨­å®š
MODEL_LLM_OPENAI = "gpt-5-nano"    # gpt-5.2, gpt-5-nano
MODEL_WHISPER_OPENAI = "whisper-1"

# OpenRouterç”¨è¨­å®š
MODEL_LLM_OPENROUTER = "nvidia/nemotron-3-nano-30b-a3b:free"    # 
MODEL_LLM_VL_OPENROUTER = "nvidia/nemotron-nano-12b-v2-vl:free" # 
OPENROUTER_BASE_URL = "https://openrouter.ai/api/v1"

# Ollamaç”¨è¨­å®š
MODEL_LLM_OLLAMA = "qwen3.5:cloud"
MODEL_LLM_VL_OLLAMA = "qwen3.5:cloud"
OLLAMA_BASE_URL = "http://localhost:11434/v1"

# å…±é€šè¨­å®š
MODEL_WHISPER = "small"     # tiny, base, small, medium, large
_WHISPER_MODEL = None


# ==== å…±é€šé–¢æ•°ç¾¤ ====

def get_whisper_model():
    """faster-whisperãƒ¢ãƒ‡ãƒ«ã‚’é…å»¶èª­ã¿è¾¼ã¿"""
    global _WHISPER_MODEL
    if _WHISPER_MODEL is None:
        from faster_whisper import WhisperModel
        _WHISPER_MODEL = WhisperModel(MODEL_WHISPER, device="auto", compute_type="int8")
    return _WHISPER_MODEL


def get_model_config(provider):
    """ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã”ã¨ã®ãƒ¢ãƒ‡ãƒ«è¨­å®šã‚’å–å¾—"""
    if provider == LLMProvider.OPENAI:
        return {
            "llm": MODEL_LLM_OPENAI,
            "whisper": MODEL_WHISPER_OPENAI,
            "vl": MODEL_LLM_OPENAI
        }
    elif provider == LLMProvider.OPENROUTER:
        return {
            "llm": MODEL_LLM_OPENROUTER,
            "whisper": MODEL_WHISPER,
            "vl": MODEL_LLM_VL_OPENROUTER
        }
    elif provider == LLMProvider.OLLAMA:
        return {
            "llm": MODEL_LLM_OLLAMA,
            "whisper": MODEL_WHISPER,
            "vl": MODEL_LLM_VL_OLLAMA
        }


def create_client(provider, api_key=None):
    """LLMãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã”ã¨ã®ã‚¯ãƒ©ã‚¤ã‚¢ãƒ³ãƒˆã‚’ä½œæˆ"""
    if provider == LLMProvider.OPENAI:
        if not api_key:
            api_key = os.environ.get('OPENAI_API_KEY')
        return openai.OpenAI(api_key=api_key)
    
    elif provider == LLMProvider.OPENROUTER:
        if not api_key:
            api_key = os.environ.get('OPENROUTER_API_KEY')
        extra_headers = {}
        site_url = os.environ.get("OPENROUTER_SITE_URL")
        app_name = os.environ.get("OPENROUTER_APP_NAME")
        if site_url:
            extra_headers["HTTP-Referer"] = site_url
        if app_name:
            extra_headers["X-Title"] = app_name

        client_kwargs = {
            "api_key": api_key,
            "base_url": OPENROUTER_BASE_URL
        }
        if extra_headers:
            client_kwargs["default_headers"] = extra_headers

        return openai.OpenAI(**client_kwargs)
    
    elif provider == LLMProvider.OLLAMA:
        # Ollamaã¯APIã‚­ãƒ¼ã‚’å¿…è¦ã¨ã—ãªã„
        return openai.OpenAI(
            api_key="ollama",  # ãƒ€ãƒŸãƒ¼ã®APIã‚­ãƒ¼ï¼ˆOpenAIã‚¯ãƒ©ã‚¤ã‚¢ãƒ³ãƒˆã§å¿…è¦ï¼‰
            base_url=OLLAMA_BASE_URL
        )


def transcribe_audio(file_path, client, provider):
    """éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ«ã‚’ãƒ†ã‚­ã‚¹ãƒˆã«å¤‰æ›"""
    if provider == LLMProvider.OPENAI:
        # OpenAI Whisper APIã‚’ä½¿ç”¨
        audio_file = open(file_path, "rb")
        response = client.audio.transcriptions.create(
            model=MODEL_WHISPER_OPENAI,
            file=audio_file,
            response_format="verbose_json",
            language="ja"
        )
        text = response.text
        segments = response.segments
    else:
        # faster-whisperã‚’ä½¿ç”¨ (OpenRouter/Ollama)
        model = get_whisper_model()
        segments_iter, _ = model.transcribe(file_path, language="ja", vad_filter=True)
        segments = list(segments_iter)
        text = " ".join(seg.text.strip() for seg in segments).strip()

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"transcription_{timestamp}.txt"
    with open(filename, "w", encoding="utf-8") as f:
        f.write(text)

    return text, segments


def analyze_speech(segments):
    """éŸ³å£°åˆ†æ: WPMã€ãƒ•ã‚£ãƒ©ãƒ¼ãƒ¯ãƒ¼ãƒ‰ã€é•·ã„é–“ã®æ¤œå‡º"""
    if not segments:
        return {
            "wpm": 0,
            "filler_count": 0,
            "long_pauses": 0
        }
    
    total_words = sum(len(seg.text.split()) for seg in segments)
    duration_minutes = (segments[-1].end - segments[0].start) / 60.0
    wpm = total_words / duration_minutes if duration_minutes else 0

    filler_words = ['ãˆãƒ¼ã¨', 'ã‚ã®', 'ãˆã£ã¨', 'ãã®']
    filler_count = sum(sum(word in seg.text for word in filler_words) for seg in segments)

    pause_lengths = [segments[i + 1].start - segments[i].end for i in range(len(segments) - 1)]
    long_pauses = sum(1 for p in pause_lengths if p > 1.0)

    return {
        "wpm": round(wpm, 2),
        "filler_count": filler_count,
        "long_pauses": long_pauses
    }


def extract_ppt_text(file_path):
    """PowerPointã‹ã‚‰ãƒ†ã‚­ã‚¹ãƒˆã‚’æŠ½å‡º"""
    prs = Presentation(file_path)
    slides_text = []
    for i, slide in enumerate(prs.slides):
        slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        slides_text.append(f"ã‚¹ãƒ©ã‚¤ãƒ‰ {i + 1}:\n{slide_text}\n")
    return "\n".join(slides_text)


def extract_images_from_ppt(ppt_path, output_dir):
    """PowerPointã‹ã‚‰ç”»åƒã‚’æŠ½å‡º"""
    prs = Presentation(ppt_path)
    image_files = []

    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    for i, slide in enumerate(prs.slides):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_filename = os.path.join(output_dir, f"slide_{i + 1}_image_{j + 1}.png")
                with open(image_filename, 'wb') as f:
                    f.write(image_bytes)
                image_files.append(image_filename)

    return image_files


def encode_image_to_base64(image_path):
    """ç”»åƒã‚’Base64ã‚¨ãƒ³ã‚³ãƒ¼ãƒ‰"""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')


def analyze_image(image_path, client, model_config):
    """ç”»åƒã‚’AIã§åˆ†æ"""
    print(f"ãƒ•ã‚¡ã‚¤ãƒ«å: {image_path}")  # ç”»åƒåˆ†æã§å¤±æ•—ã™ã‚‹å¯èƒ½æ€§ãŒã‚ã‚‹ãŸã‚ã€ãƒ‡ãƒãƒƒã‚°ç”¨ã«å‡ºåŠ›ã‚’è¿½åŠ 
    base64_image = encode_image_to_base64(image_path)
    response = client.chat.completions.create(
        model=model_config["vl"],
        messages=[
            {"role": "system", "content": "ã‚ãªãŸã¯ç”»åƒè§£æã®å°‚é–€å®¶ã§ã™ã€‚"},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "ã“ã®ç”»åƒã«ä½•ãŒå†™ã£ã¦ã„ã‚‹ã‹èª¬æ˜ã—ã€ãƒ—ãƒ¬ã‚¼ãƒ³è³‡æ–™ã¨ã—ã¦é©åˆ‡ã‹è©•ä¾¡ã—ã¦ãã ã•ã„ã€‚è¦–è¦šè³‡æ–™ã¨ã—ã¦ã®è³ªã‚‚100ç‚¹æº€ç‚¹(æ•´æ•°)ã§æ¡ç‚¹ã—ã¦ãã ã•ã„ã€‚ãƒ•ã‚©ãƒ¼ãƒãƒƒãƒˆ: è¦–è¦šè³‡æ–™: â—‹ç‚¹"},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                ]
            }
        ]
    )
    return response.choices[0].message.content


def extract_visual_score(image_analysis):
    """ç”»åƒåˆ†æçµæœã‹ã‚‰ã‚¹ã‚³ã‚¢ã‚’æŠ½å‡º"""
    pattern = r"è¦–è¦šè³‡æ–™: ([0-9]{1,3})ç‚¹"
    matches = re.findall(pattern, image_analysis)
    if matches:
        scores = [int(score) for score in matches]
        return int(sum(scores) / len(scores))
    return 0


def analyze_all_images(image_files, client, model_config):
    """å…¨ç”»åƒã‚’åˆ†æ"""
    all_analyses = []
    for image_path in image_files:
        analysis = analyze_image(image_path, client, model_config)
        all_analyses.append(f"{image_path}:\n{analysis}\n")
    return "\n".join(all_analyses)


def analyze_slide_text(slide_text, client, model_config):
    """ã‚¹ãƒ©ã‚¤ãƒ‰ãƒ†ã‚­ã‚¹ãƒˆã‚’åˆ†æ"""
    response = client.chat.completions.create(
        model=model_config["llm"],
        messages=[
            {"role": "system", "content": "ã‚ãªãŸã¯ãƒ—ãƒ­ã®ãƒ—ãƒ¬ã‚¼ãƒ³è³‡æ–™è©•ä¾¡è€…ã§ã™ã€‚"},
            {"role": "user", "content": f"""
ä»¥ä¸‹ã¯ãƒ—ãƒ¬ã‚¼ãƒ³ãƒ†ãƒ¼ã‚·ãƒ§ãƒ³ã®ã‚¹ãƒ©ã‚¤ãƒ‰å…¨æ–‡ã§ã™ã€‚

[ã‚¹ãƒ©ã‚¤ãƒ‰å…¨æ–‡]
{slide_text}

ã“ã®è³‡æ–™ã®ã‚¹ãƒ©ã‚¤ãƒ‰æ•°ã€å„ã‚¹ãƒ©ã‚¤ãƒ‰ã®æ–‡å­—é‡ã®é©åˆ‡ã•ã€å†…å®¹ã‚’è©•ä¾¡ã—ã€å…¨ä½“çš„ãªè³‡æ–™ã®è³ªã‚’ä»¥ä¸‹ã®ãƒ•ã‚©ãƒ¼ãƒãƒƒãƒˆã§100ç‚¹æº€ç‚¹(æ•´æ•°)ã§è©•ä¾¡ã—ã¦ãã ã•ã„ã€‚
ãŸã ã—ã€å›³è¡¨ã‚„ç”»åƒã¯è©•ä¾¡ã«å«ã‚ãªã„ã§ãã ã•ã„ã€‚

è³‡æ–™: â—‹ç‚¹

ãã®å¾Œã«è³‡æ–™ã®è‰¯ã„ç‚¹ã¨æ”¹å–„ç‚¹ã‚’ç°¡å˜ã«ã¾ã¨ã‚ã¦ãã ã•ã„ã€‚
"""}
        ]
    )
    return response.choices[0].message.content


def generate_evaluation_with_images(transcription, slide_text_analysis, image_analysis, client, model_config):
    """ç·åˆè©•ä¾¡ã‚’ç”Ÿæˆ"""
    response = client.chat.completions.create(
        model=model_config["llm"],
        messages=[
            {"role": "system", "content": "ã‚ãªãŸã¯ãƒ—ãƒ­ã®ãƒ—ãƒ¬ã‚¼ãƒ³è©•ä¾¡è€…ã§ã™ã€‚"},
            {"role": "user", "content": f"""
ä»¥ä¸‹ã¯ãƒ—ãƒ¬ã‚¼ãƒ³ã®æ–‡å­—èµ·ã“ã—ã¨ã‚¹ãƒ©ã‚¤ãƒ‰è³‡æ–™ã®åˆ†æçµæœã€ãŠã‚ˆã³ç”»åƒåˆ†æçµæœã§ã™ã€‚

[æ–‡å­—èµ·ã“ã—]:
{transcription}

[ã‚¹ãƒ©ã‚¤ãƒ‰ãƒ†ã‚­ã‚¹ãƒˆåˆ†æ]:
{slide_text_analysis}

[ç”»åƒåˆ†æ]:
{image_analysis}

ä»¥ä¸‹ã®4ã¤ã®è¦³ç‚¹(å†…å®¹ã€ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“ã€è¦–è¦šè³‡æ–™ã€æ§‹æˆ)ã«ã¤ã„ã¦ã€ãã‚Œãã‚Œ100ç‚¹æº€ç‚¹(æ•´æ•°)ã§è©•ä¾¡ã—ã€ç°¡å˜ãªç†ç”±ã¨æ”¹å–„ç‚¹ã€é•·æ‰€ã‚’å‡ºåŠ›ã—ã¦ãã ã•ã„ã€‚
æœ€å¾Œã«3ã¤ã®æ”¹å–„ç‚¹ã¨å…·ä½“çš„ãªã‚¢ãƒ‰ãƒã‚¤ã‚¹ã‚‚ç¤ºã—ã¦ãã ã•ã„ã€‚

ãƒ•ã‚©ãƒ¼ãƒãƒƒãƒˆã¯å¿…ãšä»¥ä¸‹ã¨ã—ã¦ãã ã•ã„:
å†…å®¹: â—‹ç‚¹
ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“: â—‹ç‚¹
è¦–è¦šè³‡æ–™: â—‹ç‚¹
æ§‹æˆ: â—‹ç‚¹

ãã®å¾Œã«è©•ä¾¡ã‚³ãƒ¡ãƒ³ãƒˆã‚’æ›¸ã„ã¦ãã ã•ã„ã€‚
"""}
        ]
    )
    return response.choices[0].message.content


def extract_scores(evaluation_text):
    """è©•ä¾¡ãƒ†ã‚­ã‚¹ãƒˆã‹ã‚‰ã‚¹ã‚³ã‚¢ã‚’æŠ½å‡º"""
    pattern = r"å†…å®¹: ([0-9]{1,3})ç‚¹.*?ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“: ([0-9]{1,3})ç‚¹.*?è¦–è¦šè³‡æ–™: ([0-9]{1,3})ç‚¹.*?æ§‹æˆ: ([0-9]{1,3})ç‚¹"
    match = re.search(pattern, evaluation_text, re.DOTALL)

    if match:
        return {
            "å†…å®¹": int(match.group(1)),
            "ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“": int(match.group(2)),
            "è¦–è¦šè³‡æ–™": int(match.group(3)),
            "æ§‹æˆ": int(match.group(4))
        }
    else:
        print("ã‚¹ã‚³ã‚¢ã®æŠ½å‡ºã«å¤±æ•—ã—ã¾ã—ãŸã€‚ãƒ‡ãƒ•ã‚©ãƒ«ãƒˆã§å…¨ã¦0ç‚¹ã¨ã—ã¾ã™ã€‚")
        return {
            "å†…å®¹": 0,
            "ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“": 0,
            "è¦–è¦šè³‡æ–™": 0,
            "æ§‹æˆ": 0
        }


def compute_score(sub_scores):
    """ã‚µãƒ–ã‚¹ã‚³ã‚¢ã‹ã‚‰ç·åˆã‚¹ã‚³ã‚¢ã‚’è¨ˆç®—"""
    weights = {
        "å†…å®¹": 0.3,
        "ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“": 0.3,
        "è¦–è¦šè³‡æ–™": 0.2,
        "æ§‹æˆ": 0.2
    }
    total = sum(float(sub_scores[k]) * weights[k] for k in weights)
    return int(round(total, 0))


def evaluate_presentation_core(audio_path, ppt_path, client, provider, progress_callback=None):
    """
    ãƒ—ãƒ¬ã‚¼ãƒ³è©•ä¾¡ã®ã‚³ã‚¢å‡¦ç†
    progress_callback: é€²æ—ã‚’é€šçŸ¥ã™ã‚‹ã‚³ãƒ¼ãƒ«ãƒãƒƒã‚¯é–¢æ•°(GUIç”¨)
    """
    model_config = get_model_config(provider)
    
    def log(message):
        print(message)
        if progress_callback:
            progress_callback(message)

    log("éŸ³å£°åˆ†æä¸­")
    text, segments = transcribe_audio(audio_path, client, provider)
    speech_analysis = analyze_speech(segments)

    log("è³‡æ–™åˆ†æä¸­")
    slides_text = extract_ppt_text(ppt_path)
    slide_text_analysis = analyze_slide_text(slides_text, client, model_config)

    log("ç”»åƒè§£æä¸­")
    image_files = extract_images_from_ppt(ppt_path, "extracted_images")
    if image_files:
        image_analysis = analyze_all_images(image_files, client, model_config)
        image_visual_score = extract_visual_score(image_analysis)
    else:
        image_analysis = "ç”»åƒã¯å«ã¾ã‚Œã¦ã„ã¾ã›ã‚“ã€‚"
        image_visual_score = 0

    log("ç·åˆè©•ä¾¡ç”Ÿæˆä¸­")
    evaluation = generate_evaluation_with_images(text, slide_text_analysis, image_analysis, client, model_config)

    sub_scores = extract_scores(evaluation)
    # ç”»åƒãªã—ã®å ´åˆã¯ç”»åƒã®å¾—ç‚¹ã‚’åˆ¤å®šã—ãªã„ã‚ˆã†ã«ã™ã‚‹
    if image_visual_score == 0:
        image_visual_score = sub_scores["è¦–è¦šè³‡æ–™"]
    sub_scores["è¦–è¦šè³‡æ–™"] = int(round((sub_scores["è¦–è¦šè³‡æ–™"] + image_visual_score) / 2, 0))

    total_score = compute_score(sub_scores)

    # çµæœã‚’ä¿å­˜
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    result_filename = f"evaluation_result_{timestamp}.txt"

    with open(result_filename, "w", encoding="utf-8") as f:
        f.write(f"==== ç·åˆå¾—ç‚¹: {total_score}ç‚¹ ====\n\n")
        f.write("==== ç·åˆè©•ä¾¡ ====\n")
        f.write(evaluation + "\n\n")
        f.write("==== éŸ³å£°åˆ†æ ====\n")
        f.write(str(speech_analysis) + "\n\n")
        f.write("==== ä½¿ç”¨ãƒ¢ãƒ‡ãƒ« ====\n")
        f.write(f"- ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼: {provider.value}\n")
        f.write(f"- LLM(å†…å®¹): {model_config['llm']}\n")
        f.write(f"- LLM(ç”»åƒ): {model_config['vl']}\n")
        f.write(f"- éŸ³å£°: {model_config['whisper']}\n")

    log(f"è©•ä¾¡çµæœã‚’ãƒ•ã‚¡ã‚¤ãƒ«ã«ä¿å­˜ã—ã¾ã—ãŸ: {result_filename}")

    # ä¸€æ™‚ç”»åƒãƒ•ã‚¡ã‚¤ãƒ«ã‚’å‰Šé™¤
    if image_files:
        for image_path in image_files:
            if os.path.exists(image_path):
                os.remove(image_path)
        if os.path.exists("extracted_images"):
            os.rmdir("extracted_images")
        log("ä¸€æ™‚ç”»åƒãƒ•ã‚¡ã‚¤ãƒ«ã‚’å‰Šé™¤ã—ã¾ã—ãŸã€‚")

    return {
        "total_score": total_score,
        "sub_scores": sub_scores,
        "evaluation": evaluation,
        "speech_analysis": speech_analysis,
        "transcription": text,
        "slide_text_analysis": slide_text_analysis,
        "image_analysis": image_analysis,
        "image_files": image_files,
        "provider": provider.value,
        "model_config": model_config
    }


# ==== CLIãƒ¢ãƒ¼ãƒ‰ ====
def run_cli_mode():
    """ã‚³ãƒãƒ³ãƒ‰ãƒ©ã‚¤ãƒ³å®Ÿè¡Œãƒ¢ãƒ¼ãƒ‰"""
    # ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã®æŒ‡å®š
    provider_name = os.environ.get('LLM_PROVIDER', PROVIDER).lower()
    try:
        provider = LLMProvider(provider_name)
    except ValueError:
        print(f"ã‚¨ãƒ©ãƒ¼: ä¸æ­£ãªãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼å '{provider_name}'")
        print(f"æœ‰åŠ¹ãªãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼: {[p.value for p in LLMProvider]}")
        sys.exit(1)

    if len(sys.argv) != 3:
        print("ä½¿ç”¨æ–¹æ³•: python presentation_evaluator.py éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ« ãƒ‘ãƒ¯ãƒãƒ•ã‚¡ã‚¤ãƒ«")
        print("ä¾‹: python presentation_evaluator.py sample.wav slides.pptx")
        print(f"\nç¾åœ¨è¨­å®šã•ã‚Œã¦ã„ã‚‹ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼: {provider.value}")
        print("ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã‚’å¤‰æ›´ã™ã‚‹ã«ã¯ç’°å¢ƒå¤‰æ•° LLM_PROVIDER ã‚’è¨­å®šã—ã¦ãã ã•ã„")
        sys.exit(1)

    audio_path = sys.argv[1]
    ppt_path = sys.argv[2]

    if not os.path.exists(audio_path):
        print(f"éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ«ãŒè¦‹ã¤ã‹ã‚Šã¾ã›ã‚“: {audio_path}")
        sys.exit(1)
    if not os.path.exists(ppt_path):
        print(f"PowerPointãƒ•ã‚¡ã‚¤ãƒ«ãŒè¦‹ã¤ã‹ã‚Šã¾ã›ã‚“: {ppt_path}")
        sys.exit(1)

    # APIã‚­ãƒ¼å–å¾—
    api_key = None
    if provider == LLMProvider.OPENAI:
        api_key = os.environ.get('OPENAI_API_KEY')
        if not api_key:
            print("ã‚¨ãƒ©ãƒ¼: OPENAI_API_KEYç’°å¢ƒå¤‰æ•°ãŒè¨­å®šã•ã‚Œã¦ã„ã¾ã›ã‚“")
            sys.exit(1)
    elif provider == LLMProvider.OPENROUTER:
        api_key = os.environ.get('OPENROUTER_API_KEY')
        if not api_key:
            print("ã‚¨ãƒ©ãƒ¼: OPENROUTER_API_KEYç’°å¢ƒå¤‰æ•°ãŒè¨­å®šã•ã‚Œã¦ã„ã¾ã›ã‚“")
            sys.exit(1)
    # Ollamaã¯APIã‚­ãƒ¼ä¸è¦

    client = create_client(provider, api_key)
    model_config = get_model_config(provider)
    
    print(f"ä½¿ç”¨ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼: {provider.value}")
    print(f"ä½¿ç”¨LLM: {model_config['llm']}")
    
    evaluate_presentation_core(audio_path, ppt_path, client, provider)


# ==== GUIãƒ¢ãƒ¼ãƒ‰ ====
def run_gui_mode():
    """Streamlit GUIãƒ¢ãƒ¼ãƒ‰"""
    if not STREAMLIT_AVAILABLE:
        print("ã‚¨ãƒ©ãƒ¼: StreamlitãŒã‚¤ãƒ³ã‚¹ãƒˆãƒ¼ãƒ«ã•ã‚Œã¦ã„ã¾ã›ã‚“")
        print("ã‚¤ãƒ³ã‚¹ãƒˆãƒ¼ãƒ«: pip install streamlit")
        sys.exit(1)

    # ãƒšãƒ¼ã‚¸è¨­å®š
    st.set_page_config(page_title="AIãƒ—ãƒ¬ã‚¼ãƒ³è©•ä¾¡ã‚·ã‚¹ãƒ†ãƒ ", layout="wide")

    # ã‚¹ã‚¿ã‚¤ãƒ«è¨­å®š
    st.markdown("""
        <style>
        .main {
            background-color: #f8f9fa;
        }
        .stMetric {
            background-color: #ffffff;
            padding: 15px;
            border-radius: 10px;
            box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        }
        </style>
        """, unsafe_allow_html=True)

    # ã‚¿ã‚¤ãƒˆãƒ«ãƒ»èª¬æ˜
    st.title("ğŸ¤ AIãƒ—ãƒ¬ã‚¼ãƒ³è©•ä¾¡ã‚·ã‚¹ãƒ†ãƒ ")
    st.markdown("éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ«ã¨ã‚¹ãƒ©ã‚¤ãƒ‰è³‡æ–™ã‚’ã‚¢ãƒƒãƒ—ãƒ­ãƒ¼ãƒ‰ã™ã‚‹ã ã‘ã§ã€AIãŒã‚ãªãŸã®ãƒ—ãƒ¬ã‚¼ãƒ³ã‚’å¤šè§’çš„ã«åˆ†æãƒ»æ¡ç‚¹ã—ã¾ã™ã€‚")

    # ã‚µã‚¤ãƒ‰ãƒãƒ¼è¨­å®š
    with st.sidebar:
        st.header("âš™ï¸ è¨­å®š")
        
        # ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼é¸æŠ
        provider_select = st.selectbox(
            "LLMãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã‚’é¸æŠ",
            options=[p.value for p in LLMProvider],
            index=[p.value for p in LLMProvider].index(PROVIDER)
        )
        provider = LLMProvider(provider_select)
        model_config = get_model_config(provider)
        
        # APIã‚­ãƒ¼å…¥åŠ›ï¼ˆãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã«ã‚ˆã£ã¦è¡¨ç¤ºåˆ‡æ›¿ï¼‰
        api_key = None
        if provider == LLMProvider.OPENAI:
            api_key = st.text_input("OpenAI API Keyã‚’å…¥åŠ›ã—ã¦ãã ã•ã„", type="password")
            if not api_key:
                st.warning("âš ï¸ ç¶šè¡Œã™ã‚‹ã«ã¯OpenAI APIã‚­ãƒ¼ã‚’å…¥åŠ›ã—ã¦ãã ã•ã„ã€‚")
                st.stop()
        elif provider == LLMProvider.OPENROUTER:
            api_key = st.text_input("OpenRouter API Keyã‚’å…¥åŠ›ã—ã¦ãã ã•ã„", type="password")
            if not api_key:
                st.warning("âš ï¸ ç¶šè¡Œã™ã‚‹ã«ã¯OpenRouter APIã‚­ãƒ¼ã‚’å…¥åŠ›ã—ã¦ãã ã•ã„ã€‚")
                st.stop()
        else:
            st.info("Ollamaã‚’ä½¿ç”¨ã—ã¾ã™ã€‚OllamaãŒèµ·å‹•ã—ã¦ã„ã‚‹ã“ã¨ã‚’ç¢ºèªã—ã¦ãã ã•ã„ã€‚")
        
        st.info(f"""
        **ä½¿ç”¨ãƒ¢ãƒ‡ãƒ«:**
        - LLM(å†…å®¹): {model_config['llm']}
        - LLM(ç”»åƒ): {model_config['vl']}
        - éŸ³å£°: {model_config['whisper']}
        
        **åˆ†æé …ç›®:**
        1. å†…å®¹ (30%)
        2. ãƒ—ãƒ¬ã‚¼ãƒ³æŠ€è¡“ (30%)
        3. è¦–è¦šè³‡æ–™ (20%)
        4. æ§‹æˆ (20%)
        """)

    client = create_client(provider, api_key)

    # ãƒ•ã‚¡ã‚¤ãƒ«ã‚¢ãƒƒãƒ—ãƒ­ãƒ¼ãƒ‰
    col1, col2 = st.columns(2)
    with col1:
        audio_upload = st.file_uploader("1. éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ«ã‚’ã‚¢ãƒƒãƒ—ãƒ­ãƒ¼ãƒ‰", type=['mp3', 'wav', 'm4a', 'mp4'])
    with col2:
        ppt_upload = st.file_uploader("2. PowerPointãƒ•ã‚¡ã‚¤ãƒ«ã‚’ã‚¢ãƒƒãƒ—ãƒ­ãƒ¼ãƒ‰", type=['pptx'])

    if st.button("ğŸ“Š ãƒ—ãƒ¬ã‚¼ãƒ³ã‚’åˆ†æã™ã‚‹", use_container_width=True):
        if audio_upload and ppt_upload:
            # ä¸€æ™‚ãƒ•ã‚¡ã‚¤ãƒ«ä¿å­˜
            temp_dir = "temp_process"
            if not os.path.exists(temp_dir):
                os.makedirs(temp_dir)
            
            audio_path = os.path.join(temp_dir, audio_upload.name)
            ppt_path = os.path.join(temp_dir, ppt_upload.name)
            
            with open(audio_path, "wb") as f:
                f.write(audio_upload.getbuffer())
            with open(ppt_path, "wb") as f:
                f.write(ppt_upload.getbuffer())

            try:
                with st.status("åˆ†æä¸­...", expanded=True) as status:
                    progress_messages = []
                    
                    def progress_callback(msg):
                        progress_messages.append(msg)
                        icon_map = {
                            "éŸ³å£°åˆ†æä¸­": "ğŸ™ï¸",
                            "è³‡æ–™åˆ†æä¸­": "ğŸ“„",
                            "ç”»åƒè§£æä¸­": "ğŸ–¼ï¸",
                            "ç·åˆè©•ä¾¡ç”Ÿæˆä¸­": "ğŸ¤–"
                        }
                        icon = icon_map.get(msg, "â³")
                        st.write(f"{icon} {msg}...")

                    result = evaluate_presentation_core(
                        audio_path, ppt_path, client, provider,
                        progress_callback=progress_callback
                    )

                    status.update(label="âœ… åˆ†æãŒå®Œäº†ã—ã¾ã—ãŸï¼", state="complete", expanded=False)

                # çµæœè¡¨ç¤º
                st.divider()
                
                tab1, tab2, tab3 = st.tabs(["ğŸ“ ç·åˆè©•ä¾¡ãƒ¬ãƒãƒ¼ãƒˆ", "ğŸ“– æ–‡å­—èµ·ã“ã—å…¨æ–‡", "ğŸ–¼ï¸ ã‚¹ãƒ©ã‚¤ãƒ‰åˆ†æè©³ç´°"])
                
                with tab1:
                    st.subheader(f"ğŸ“Š ç·åˆã‚¹ã‚³ã‚¢: {result['total_score']} ç‚¹")
                    
                    cols = st.columns(4)
                    for i, (label, score) in enumerate(result['sub_scores'].items()):
                        cols[i].caption(f"{label}: {score}ç‚¹")
                    
                    st.markdown("---")
                    st.markdown(result['evaluation'])
                    
                with tab2:
                    st.text_area("æ–‡å­—èµ·ã“ã—å†…å®¹", result['transcription'], height=300)
                    
                with tab3:
                    st.markdown("### ã‚¹ãƒ©ã‚¤ãƒ‰ãƒ†ã‚­ã‚¹ãƒˆè©•ä¾¡")
                    st.write(result['slide_text_analysis'])
                    if result['image_files']:
                        st.markdown("### æŠ½å‡ºã•ã‚ŒãŸç”»åƒã¨AIã‚³ãƒ¡ãƒ³ãƒˆ")
                        for img in result['image_files']:
                            if os.path.exists(img):
                                st.image(img, width=300)

            except Exception as e:
                st.error(f"åˆ†æä¸­ã«ã‚¨ãƒ©ãƒ¼ãŒç™ºç”Ÿã—ã¾ã—ãŸ: {e}")
            
            finally:
                # ã‚¯ãƒªãƒ¼ãƒ³ã‚¢ãƒƒãƒ—
                if os.path.exists(temp_dir):
                    shutil.rmtree(temp_dir)

        else:
            st.info("éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ«ã¨PowerPointãƒ•ã‚¡ã‚¤ãƒ«ã‚’ä¸¡æ–¹ã‚¢ãƒƒãƒ—ãƒ­ãƒ¼ãƒ‰ã—ã¦ã€åˆ†æãƒœã‚¿ãƒ³ã‚’æŠ¼ã—ã¦ãã ã•ã„ã€‚")

    # ãƒ•ãƒƒã‚¿ãƒ¼
    st.markdown("---")
    st.caption(
        f"Presentation Evaluator Pro v2.0 (çµ±åˆç‰ˆ) | Powered by {provider.value.upper()} | LLM: {model_config['llm']}"
    )


# ==== ãƒ¡ã‚¤ãƒ³å®Ÿè¡Œéƒ¨ ====
if __name__ == "__main__":
    # ã‚³ãƒãƒ³ãƒ‰ãƒ©ã‚¤ãƒ³å¼•æ•°ãŒã‚ã‚Œã°CLIãƒ¢ãƒ¼ãƒ‰ã€ãªã‘ã‚Œã°GUIãƒ¢ãƒ¼ãƒ‰
    if len(sys.argv) > 1:
        # CLIãƒ¢ãƒ¼ãƒ‰ã§å®Ÿè¡Œ
        run_cli_mode()
    else:
        # GUIãƒ¢ãƒ¼ãƒ‰ã§å®Ÿè¡Œï¼ˆStreamlitã‹ã‚‰èµ·å‹•ã•ã‚Œã‚‹æƒ³å®šï¼‰
        if STREAMLIT_AVAILABLE:
            run_gui_mode()
        else:
            print("ã‚¨ãƒ©ãƒ¼: StreamlitãŒã‚¤ãƒ³ã‚¹ãƒˆãƒ¼ãƒ«ã•ã‚Œã¦ã„ã¾ã›ã‚“")
            print("GUIãƒ¢ãƒ¼ãƒ‰ã‚’ä½¿ç”¨ã™ã‚‹ã«ã¯: pip install streamlit")
            print("  streamlit run presentation_evaluator.py")
            print("\nCLIãƒ¢ãƒ¼ãƒ‰ã§ä½¿ç”¨ã™ã‚‹å ´åˆ:")
            print("  python presentation_evaluator.py <éŸ³å£°ãƒ•ã‚¡ã‚¤ãƒ«> <PowerPointãƒ•ã‚¡ã‚¤ãƒ«>")
            print(f"\nç’°å¢ƒå¤‰æ•° LLM_PROVIDER ã§ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼ã‚’æŒ‡å®šã§ãã¾ã™: openai, openrouter, ollama")
            print(f"ç¾åœ¨ã®ãƒ—ãƒ­ãƒã‚¤ãƒ€ãƒ¼: {PROVIDER}")
            sys.exit(1)
