import os
import json
import re
from typing import Dict, List, Tuple, Optional
from dataclasses import dataclass
from datetime import datetime
import logging

# 音声処理用ライブラリ
try:
    import whisper
    import librosa
    import numpy as np
    AUDIO_AVAILABLE = True
except ImportError:
    AUDIO_AVAILABLE = False
    print("音声処理ライブラリが不足しています。pip install whisper librosa numpy を実行してください。")

# PowerPoint処理用ライブラリ
try:
    from pptx import Presentation
    import pptx.util
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False
    print("PowerPoint処理ライブラリが不足しています。pip install python-pptx を実行してください。")

# テキスト分析用ライブラリ (spaCy)
try:
    import spacy
    from spacytextblob.spacytextblob import SpacyTextBlob
    SPACY_AVAILABLE = True
except ImportError:
    SPACY_AVAILABLE = False
    print("spaCy が不足しています。pip install spacy spacytextblob を実行し、python -m spacy download ja_core_news_sm も実行してください。")

@dataclass
class AudioAnalysis:
    """音声分析結果のデータクラス"""
    transcript: str
    duration: float
    words_per_minute: float
    filler_words_count: int
    sentiment_score: float
    volume_variation: float
    pause_count: int
    text_complexity: float  # 新規追加: テキストの複雑さ
    named_entities: List[str]  # 新規追加: 固有表現

@dataclass
class SlideAnalysis:
    """スライド分析結果のデータクラス"""
    total_slides: int
    avg_words_per_slide: float
    total_words: int
    has_images: bool
    font_consistency: float
    structure_score: float

@dataclass
class EvaluationResult:
    """評価結果のデータクラス"""
    content_score: float
    technique_score: float
    visual_score: float
    structure_score: float
    total_score: float
    strengths: List[str]
    improvements: List[str]
    detailed_advice: Dict[str, str]

class PresentationEvaluator:
    """プレゼンテーション評価システムのメインクラス"""
    
    def __init__(self):
        """初期化"""
        self.logger = self._setup_logger()
        
        # 日本語フィラーワードリスト
        self.filler_words = ['えーと', 'あの', 'その', 'えー', 'あー', 'まあ', 'なんか', 'ちょっと']
        
        # 音声処理モデルの初期化
        if AUDIO_AVAILABLE:
            try:
                self.whisper_model = whisper.load_model("base")
            except Exception as e:
                self.logger.warning(f"Whisperモデルの読み込みに失敗: {e}")
                self.whisper_model = None
        
        # spaCyの初期化
        if SPACY_AVAILABLE:
            try:
                # 日本語モデルの読み込み
                self.nlp = spacy.load("ja_core_news_sm")
                
                # TextBlobの追加（感情分析用）
                self.nlp.add_pipe('spacytextblob')
                
                self.logger.info("spaCyモデルが正常に読み込まれました")
            except OSError as e:
                self.logger.warning(f"日本語モデルの読み込みに失敗: {e}")
                self.logger.warning("python -m spacy download ja_core_news_sm を実行してください")
                self.nlp = None
            except Exception as e:
                self.logger.warning(f"spaCyの初期化に失敗: {e}")
                self.nlp = None
        else:
            self.nlp = None
    
    def _setup_logger(self) -> logging.Logger:
        """ログの設定"""
        logger = logging.getLogger('PresentationEvaluator')
        logger.setLevel(logging.INFO)
        
        if not logger.handlers:
            handler = logging.StreamHandler()
            formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
            handler.setFormatter(formatter)
            logger.addHandler(handler)
        
        return logger
    
    def _analyze_text_with_spacy(self, text: str) -> Tuple[float, float, List[str]]:
        """spaCyを使ったテキスト分析"""
        if not self.nlp or not text.strip():
            return 0.0, 0.0, []
        
        try:
            doc = self.nlp(text)
            
            # 感情スコアの計算（TextBlob使用）
            sentiment_score = doc._.blob.polarity
            
            # テキストの複雑さ計算（語彙の豊富さ、文の長さ等）
            sentences = list(doc.sents)
            tokens = [token for token in doc if not token.is_space and not token.is_punct]
            
            # 語彙の多様性（異なる語の数 / 全語数）
            unique_lemmas = set([token.lemma_ for token in tokens if not token.is_stop])
            vocab_diversity = len(unique_lemmas) / len(tokens) if tokens else 0
            
            # 平均文長
            avg_sentence_length = len(tokens) / len(sentences) if sentences else 0
            
            # 複雑さスコア（0-1の範囲）
            complexity = min(1.0, (vocab_diversity * 0.7 + min(avg_sentence_length / 20, 1.0) * 0.3))
            
            # 固有表現の抽出
            entities = [ent.text for ent in doc.ents if ent.label_ in ['PERSON', 'ORG', 'GPE', 'EVENT', 'PRODUCT']]
            
            return sentiment_score, complexity, entities
            
        except Exception as e:
            self.logger.warning(f"spaCyテキスト分析でエラー: {e}")
            return 0.0, 0.0, []
    
    def analyze_audio(self, audio_file_path: str) -> Optional[AudioAnalysis]:
        """音声ファイルの分析"""
        try:
            self.logger.info(f"音声ファイルを分析中: {audio_file_path}")
            
            if not AUDIO_AVAILABLE or not self.whisper_model:
                self.logger.error("音声処理ライブラリが利用できません")
                return None
            
            # 音声の文字起こし
            result = self.whisper_model.transcribe(audio_file_path, fp16=False, language='ja')
            transcript = result['text']
            
            # 音声データの読み込み（エラーハンドリング強化）
            try:
                y, sr = librosa.load(audio_file_path)
                duration = librosa.get_duration(y=y, sr=sr)
            except Exception as load_error:
                self.logger.warning(f"librosaでの読み込みに失敗: {load_error}")
                # Whisperの結果から時間を取得（代替手段）
                if 'segments' in result and result['segments']:
                    duration = result['segments'][-1]['end']
                else:
                    # デフォルト値を使用
                    duration = 300.0  # 5分と仮定
                    self.logger.warning("音声の長さを推定値で設定しました")
                
                # 音声分析をスキップして、基本的な分析のみ実行
                y, sr = None, None
            
            # 話速の計算（WPM）
            word_count = len(transcript.split())
            words_per_minute = (word_count / duration) * 60 if duration > 0 else 0
            
            # フィラーワードのカウント
            filler_count = sum(transcript.lower().count(word) for word in self.filler_words)
            
            # spaCyによるテキスト分析
            sentiment_score, text_complexity, entities = self._analyze_text_with_spacy(transcript)
            
            # 音量変化の分析（音声データが利用可能な場合のみ）
            volume_variation = 0.1  # デフォルト値
            pause_count = 5  # デフォルト値
            
            if y is not None and sr is not None:
                try:
                    # 音量変化の分析
                    rms = librosa.feature.rms(y=y)[0]
                    volume_variation = np.std(rms)
                    
                    # 間の分析（簡易版）
                    pause_threshold = 0.1
                    silent_samples = np.where(np.abs(y) < pause_threshold)[0]
                    pause_count = len(np.split(silent_samples, np.where(np.diff(silent_samples) != 1)[0] + 1))
                except Exception as analysis_error:
                    self.logger.warning(f"音響分析でエラー: {analysis_error}")
                    # デフォルト値を使用
            else:
                self.logger.info("音響分析はスキップし、テキスト分析のみ実行しました")
            
            return AudioAnalysis(
                transcript=transcript,
                duration=duration,
                words_per_minute=words_per_minute,
                filler_words_count=filler_count,
                sentiment_score=sentiment_score,
                volume_variation=volume_variation,
                pause_count=pause_count,
                text_complexity=text_complexity,
                named_entities=entities
            )
            
        except Exception as e:
            self.logger.error(f"音声分析でエラーが発生: {e}")
            return None
    
    def analyze_slides(self, ppt_file_path: str) -> Optional[SlideAnalysis]:
        """PowerPointファイルの分析"""
        try:
            self.logger.info(f"PowerPointファイルを分析中: {ppt_file_path}")
            
            if not PPT_AVAILABLE:
                self.logger.error("PowerPoint処理ライブラリが利用できません")
                return None
            
            prs = Presentation(ppt_file_path)
            
            total_slides = len(prs.slides)
            total_words = 0
            has_images = False
            font_sizes = []
            
            for slide in prs.slides:
                slide_text = ""
                
                # テキスト抽出
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + " "
                    
                    # 画像の存在確認
                    if shape.shape_type == 13:  # Picture type
                        has_images = True
                    
                    # フォントサイズの取得
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size:
                                    font_sizes.append(run.font.size.pt)
                
                # spaCyによる語数カウント（より正確）
                if self.nlp and slide_text.strip():
                    doc = self.nlp(slide_text)
                    words_in_slide = len([token for token in doc if not token.is_space and not token.is_punct])
                else:
                    words_in_slide = len(slide_text.split())
                
                total_words += words_in_slide
            
            avg_words_per_slide = total_words / total_slides if total_slides > 0 else 0
            
            # フォント一貫性の計算
            font_consistency = 1.0
            if font_sizes:
                font_variation = np.std(font_sizes) / np.mean(font_sizes) if np.mean(font_sizes) > 0 else 0
                font_consistency = max(0, 1 - font_variation)
            
            # 構成スコアの計算（簡易版）
            structure_score = min(1.0, total_slides / 10) * 0.5 + (1 if has_images else 0) * 0.3 + font_consistency * 0.2
            
            return SlideAnalysis(
                total_slides=total_slides,
                avg_words_per_slide=avg_words_per_slide,
                total_words=total_words,
                has_images=has_images,
                font_consistency=font_consistency,
                structure_score=structure_score
            )
            
        except Exception as e:
            self.logger.error(f"スライド分析でエラーが発生: {e}")
            return None
    
    def calculate_scores(self, audio_analysis: AudioAnalysis, slide_analysis: SlideAnalysis) -> EvaluationResult:
        """評価スコアの計算"""
        
        # 内容スコア（30%）
        content_score = self._calculate_content_score(audio_analysis, slide_analysis)
        
        # プレゼン技術スコア（30%）
        technique_score = self._calculate_technique_score(audio_analysis)
        
        # 視覚資料スコア（20%）
        visual_score = self._calculate_visual_score(slide_analysis)
        
        # 構成スコア（20%）
        structure_score = self._calculate_structure_score(audio_analysis, slide_analysis)
        
        # 総合スコア（100点満点）
        total_score = (
            content_score * 0.3 + 
            technique_score * 0.3 + 
            visual_score * 0.2 + 
            structure_score * 0.2
        ) * 20  # 5点満点を100点満点に変換
        
        # 長所と改善点の特定
        strengths = self._identify_strengths(content_score, technique_score, visual_score, structure_score, audio_analysis)
        improvements = self._identify_improvements(content_score, technique_score, visual_score, structure_score, audio_analysis)
        
        # 詳細アドバイスの生成
        detailed_advice = self._generate_detailed_advice(
            content_score, technique_score, visual_score, structure_score, 
            audio_analysis, slide_analysis
        )
        
        return EvaluationResult(
            content_score=content_score,
            technique_score=technique_score,
            visual_score=visual_score,
            structure_score=structure_score,
            total_score=total_score,
            strengths=strengths,
            improvements=improvements,
            detailed_advice=detailed_advice
        )
    
    def _calculate_content_score(self, audio: AudioAnalysis, slide: SlideAnalysis) -> float:
        """内容スコアの計算（spaCy分析結果を考慮）"""
        score = 3.0  # 基準点
        
        # 文字数による評価
        if slide.total_words > 100:
            score += 0.3
        if slide.total_words > 300:
            score += 0.3
        
        # テキストの複雑さによる評価
        if audio.text_complexity > 0.5:
            score += 0.5
        elif audio.text_complexity > 0.3:
            score += 0.3
        
        # 固有表現の使用による評価
        if len(audio.named_entities) > 3:
            score += 0.4
        elif len(audio.named_entities) > 0:
            score += 0.2
        
        # 感情表現の評価
        if abs(audio.sentiment_score) > 0.1:
            score += 0.3
        
        # 話す長さの適切性
        if 300 <= audio.duration <= 900:  # 5-15分
            score += 0.3
        
        return min(5.0, max(1.0, score))
    
    def _calculate_technique_score(self, audio: AudioAnalysis) -> float:
        """プレゼン技術スコアの計算"""
        score = 3.0  # 基準点
        
        # 話速の評価（120-180 WPMが理想的）
        if 120 <= audio.words_per_minute <= 180:
            score += 1.0
        elif 100 <= audio.words_per_minute <= 200:
            score += 0.5
        
        # フィラーワードの少なさ
        filler_ratio = audio.filler_words_count / len(audio.transcript.split()) if audio.transcript else 0
        if filler_ratio < 0.02:
            score += 1.0
        elif filler_ratio < 0.05:
            score += 0.5
        
        # 音量変化（抑揚）
        if audio.volume_variation > 0.1:
            score += 0.5
        
        return min(5.0, max(1.0, score))
    
    def _calculate_visual_score(self, slide: SlideAnalysis) -> float:
        """視覚資料スコアの計算"""
        score = 3.0  # 基準点
        
        # 適切なスライド数
        if 5 <= slide.total_slides <= 20:
            score += 0.5
        
        # 1スライドあたりの文字数
        if 20 <= slide.avg_words_per_slide <= 50:
            score += 0.5
        
        # 画像の使用
        if slide.has_images:
            score += 0.5
        
        # フォントの一貫性
        score += slide.font_consistency * 0.5
        
        return min(5.0, max(1.0, score))
    
    def _calculate_structure_score(self, audio: AudioAnalysis, slide: SlideAnalysis) -> float:
        """構成スコアの計算"""
        score = 3.0  # 基準点
        
        # スライドの構成
        score += slide.structure_score
        
        # 発表時間の適切性
        if 300 <= audio.duration <= 900:  # 5-15分
            score += 0.5
        
        # 間の使い方
        pause_ratio = audio.pause_count / audio.duration if audio.duration > 0 else 0
        if 0.1 <= pause_ratio <= 0.3:
            score += 0.5
        
        return min(5.0, max(1.0, score))
    
    def _identify_strengths(self, content: float, technique: float, visual: float, structure: float, audio: AudioAnalysis) -> List[str]:
        """長所の特定（spaCy分析結果を考慮）"""
        strengths = []
        scores = {
            "内容": content,
            "プレゼン技術": technique,
            "視覚資料": visual,
            "構成": structure
        }
        
        # 高得点の項目を長所として選出
        for category, score in sorted(scores.items(), key=lambda x: x[1], reverse=True):
            if score >= 4.0:
                if category == "内容":
                    # spaCy分析結果を活用した詳細なコメント
                    entity_comment = f"（{len(audio.named_entities)}個の固有表現を効果的に使用）" if audio.named_entities else ""
                    complexity_comment = f"語彙の豊富さが{audio.text_complexity:.1f}" if audio.text_complexity > 0.5 else ""
                    strengths.append(f"内容が充実しており、聞き手に価値のある情報を提供できています{entity_comment}")
                elif category == "プレゼン技術":
                    strengths.append("話し方が明瞭で、聞き取りやすいプレゼンテーションです")
                elif category == "視覚資料":
                    strengths.append("スライドデザインが効果的で、視覚的に理解しやすい構成です")
                elif category == "構成":
                    strengths.append("論理的な構成で、話の流れが明確です")
        
        # テキストの複雑さが高い場合の特別な評価
        if audio.text_complexity > 0.6:
            strengths.append("専門的で深みのある内容を扱っており、聞き手の知識向上に貢献します")
        
        # 固有表現が豊富な場合の評価
        if len(audio.named_entities) > 5:
            strengths.append("具体的な事例や固有名詞を多く使用し、説得力のある内容です")
        
        # 最低3つの長所を確保
        while len(strengths) < 3:
            strengths.append("プレゼンテーションに真摯に取り組む姿勢が評価できます")
        
        return strengths[:3]
    
    def _identify_improvements(self, content: float, technique: float, visual: float, structure: float, audio: AudioAnalysis) -> List[str]:
        """改善点の特定（spaCy分析結果を考慮）"""
        improvements = []
        scores = {
            "内容": content,
            "プレゼン技術": technique,
            "視覚資料": visual,
            "構成": structure
        }
        
        # 低得点の項目を改善点として選出
        for category, score in sorted(scores.items(), key=lambda x: x[1]):
            if score < 3.5:
                if category == "内容":
                    if audio.text_complexity < 0.3:
                        improvements.append("内容をより具体的で魅力的にし、専門用語や事例を増やしましょう")
                    elif len(audio.named_entities) == 0:
                        improvements.append("具体的な固有名詞や事例を追加して、説得力を高めましょう")
                    else:
                        improvements.append("内容をより具体的で魅力的にする工夫が必要です")
                elif category == "プレゼン技術":
                    improvements.append("話し方のテクニックを向上させることで、より効果的な発表になります")
                elif category == "視覚資料":
                    improvements.append("スライドデザインを改善し、視覚的インパクトを高めましょう")
                elif category == "構成":
                    improvements.append("プレゼンテーションの構成を見直し、より論理的な流れを作りましょう")
        
        # テキストの複雑さが低い場合の特別なアドバイス
        if audio.text_complexity < 0.2:
            improvements.append("語彙を豊富にし、より専門的で深みのある表現を心がけましょう")
        
        # 最低3つの改善点を確保
        while len(improvements) < 3:
            improvements.append("継続的な練習により、さらなる向上が期待できます")
        
        return improvements[:3]
    
    def _generate_detailed_advice(self, content: float, technique: float, visual: float, structure: float,
                                audio: AudioAnalysis, slide: SlideAnalysis) -> Dict[str, str]:
        """詳細アドバイスの生成（spaCy分析結果を考慮）"""
        advice = {}
        
        # 内容に関するアドバイス
        if content < 4.0:
            entity_info = f"固有表現{len(audio.named_entities)}個、" if audio.named_entities else "固有表現が少なく、"
            complexity_info = f"語彙の複雑さ{audio.text_complexity:.2f}"
            advice["内容"] = f"現在の文字数は{slide.total_words}語です。{entity_info}{complexity_info}。より具体的な例やデータを追加し、専門用語を適切に使用しましょう。"
        else:
            advice["内容"] = f"内容が充実しています（語彙複雑さ: {audio.text_complexity:.2f}、固有表現: {len(audio.named_entities)}個）。この調子で、さらに深い洞察を加えていきましょう。"
        
        # プレゼン技術に関するアドバイス
        if technique < 4.0:
            wpm_advice = ""
            if audio.words_per_minute < 120:
                wpm_advice = "話速がやや遅いです。もう少しテンポを上げると良いでしょう。"
            elif audio.words_per_minute > 180:
                wpm_advice = "話速がやや速いです。もう少しゆっくり話すことを心がけましょう。"
            
            filler_advice = ""
            if audio.filler_words_count > 10:
                filler_advice = f"フィラーワードが{audio.filler_words_count}回使用されています。意識的に減らしましょう。"
            
            advice["プレゼン技術"] = f"現在の話速は{audio.words_per_minute:.1f}語/分です。{wpm_advice} {filler_advice}".strip()
        else:
            advice["プレゼン技術"] = "話し方のテクニックが優れています。この調子を維持しましょう。"
        
        # 視覚資料に関するアドバイス
        if visual < 4.0:
            advice["視覚資料"] = f"スライド数は{slide.total_slides}枚、1枚あたり平均{slide.avg_words_per_slide:.1f}語です。画像やグラフを追加し、視覚的な魅力を高めましょう。"
        else:
            advice["視覚資料"] = "視覚資料が効果的に活用されています。素晴らしいデザインセンスです。"
        
        # 構成に関するアドバイス
        if structure < 4.0:
            advice["構成"] = f"発表時間は{audio.duration/60:.1f}分です。導入・本論・結論の流れをより明確にし、論理的な構成を心がけましょう。"
        else:
            advice["構成"] = "論理的で分かりやすい構成です。聞き手にとって理解しやすいプレゼンテーションです。"
        
        return advice
    
    def evaluate_presentation(self, audio_file: str, ppt_file: str) -> Optional[EvaluationResult]:
        """プレゼンテーションの総合評価"""
        try:
            self.logger.info("プレゼンテーション評価を開始します")
            
            # ファイルの存在確認
            if not os.path.exists(audio_file):
                self.logger.error(f"音声ファイルが見つかりません: {audio_file}")
                return None
            
            if not os.path.exists(ppt_file):
                self.logger.error(f"PowerPointファイルが見つかりません: {ppt_file}")
                return None
            
            # 音声分析
            audio_analysis = self.analyze_audio(audio_file)
            if not audio_analysis:
                self.logger.error("音声分析に失敗しました")
                return None
            
            # スライド分析
            slide_analysis = self.analyze_slides(ppt_file)
            if not slide_analysis:
                self.logger.error("スライド分析に失敗しました")
                return None
            
            # 評価計算
            result = self.calculate_scores(audio_analysis, slide_analysis)
            
            self.logger.info("プレゼンテーション評価が完了しました")
            return result
            
        except Exception as e:
            self.logger.error(f"評価処理でエラーが発生: {e}")
            return None
    
    def save_result(self, result: EvaluationResult, output_file: str = None) -> str:
        """評価結果の保存"""
        if not output_file:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_file = f"presentation_evaluation_{timestamp}.json"
        
        try:
            result_dict = {
                "評価日時": datetime.now().isoformat(),
                "スコア": {
                    "内容": result.content_score,
                    "プレゼン技術": result.technique_score,
                    "視覚資料": result.visual_score,
                    "構成": result.structure_score,
                    "総合点": result.total_score
                },
                "長所": result.strengths,
                "改善点": result.improvements,
                "詳細アドバイス": result.detailed_advice
            }
            
            with open(output_file, 'w', encoding='utf-8') as f:
                json.dump(result_dict, f, ensure_ascii=False, indent=2)
            
            self.logger.info(f"評価結果を保存しました: {output_file}")
            return output_file
            
        except Exception as e:
            self.logger.error(f"結果保存でエラーが発生: {e}")
            return ""
    
    def print_result(self, result: EvaluationResult):
        """評価結果の表示"""
        print("\n" + "="*60)
        print("プレゼンテーション評価結果")
        print("="*60)
        
        print(f"\n【総合得点】 {result.total_score:.1f}/100点")
        
        print(f"\n【項目別スコア】")
        print(f"内容:         {result.content_score:.1f}/5.0")
        print(f"プレゼン技術: {result.technique_score:.1f}/5.0")
        print(f"視覚資料:     {result.visual_score:.1f}/5.0")
        print(f"構成:         {result.structure_score:.1f}/5.0")
        
        print(f"\n【長所】")
        for i, strength in enumerate(result.strengths, 1):
            print(f"{i}. {strength}")
        
        print(f"\n【改善点】")
        for i, improvement in enumerate(result.improvements, 1):
            print(f"{i}. {improvement}")
        
        print(f"\n【詳細アドバイス】")
        for category, advice in result.detailed_advice.items():
            print(f"■ {category}")
            print(f"  {advice}")
        
        print("\n" + "="*60)

def main():
    """メイン関数"""
    print("プレゼンテーション評価システム")
    print("================================")
    
    # 使用例
    evaluator = PresentationEvaluator()
    
    # ファイルパスの入力
#    audio_file = input("音声ファイルのパスを入力してください: ").strip()
#    ppt_file = input("PowerPointファイルのパスを入力してください: ").strip()
    audio_file = "C:\work\pyemv\py310\PresenEvaluator\honma_202505.mp3" # type: ignore
    ppt_file = "C:\work\pyemv\py310\PresenEvaluator\honma_202505.pptx" # type: ignore
    
    if not audio_file or not ppt_file:
        print("ファイルパスが入力されていません。")
        return
    
    # 評価実行
    result = evaluator.evaluate_presentation(audio_file, ppt_file)
    
    if result:
        # 結果表示
        evaluator.print_result(result)
        
        # 結果保存
        save_choice = input("\n結果をファイルに保存しますか？ (y/n): ").strip().lower()
        if save_choice == 'y':
            output_file = evaluator.save_result(result)
            if output_file:
                print(f"結果が保存されました: {output_file}")
    else:
        print("評価に失敗しました。ファイルパスとライブラリのインストール状況を確認してください。")

if __name__ == "__main__":
    main()